"""Format-aware document parser. Returns a list of source-tagged text chunks.

Supported formats:
  PDF    — pypdf for digital text; PyMuPDF + Gemini Vision fallback for scanned pages
  PPTX   — slide text + speaker notes
  CSV    — row-grouped chunks
  XLSX   — ALL sheets, row-grouped chunks per sheet
  DOCX   — paragraph + table content, sliding window chunking
  TXT/MD — sliding window chunking
"""
import logging
import os
from typing import List

logger = logging.getLogger(__name__)

_CHUNK_SIZE = 600
_CHUNK_OVERLAP = 80
_ROWS_PER_CHUNK = 15


def _sliding_window(text: str, size: int = _CHUNK_SIZE, overlap: int = _CHUNK_OVERLAP) -> List[str]:
    if not text:
        return []
    chunks = []
    step = max(size - overlap, 1)
    for start in range(0, len(text), step):
        piece = text[start : start + size]
        if piece.strip():
            chunks.append(piece)
        if start + size >= len(text):
            break
    return chunks


def _df_to_chunks(df, filename: str, sheet_label: str) -> List[str]:
    """Convert a DataFrame to source-tagged row-group chunks."""
    import pandas as pd

    if df.empty:
        return []

    cols = list(df.columns)
    rows_text = []
    for _, row in df.iterrows():
        parts = []
        for c in cols:
            val = row[c]
            if pd.notna(val) and str(val).strip():
                parts.append(f"{c}={val}")
        if parts:
            rows_text.append(", ".join(parts))

    chunks = []
    for i in range(0, len(rows_text), _ROWS_PER_CHUNK):
        batch = rows_text[i : i + _ROWS_PER_CHUNK]
        row_start = i + 1
        row_end = i + len(batch)
        block = "\n".join(batch)
        tagged = f"[source: {filename} | {sheet_label} rows {row_start}-{row_end}]\n{block}"
        chunks.append(tagged)
    return chunks


def _gemini_ocr_page(img_bytes: bytes, filename: str, page_num: int) -> str:
    """Send a rendered page image to Gemini Vision and return extracted text."""
    try:
        import google.generativeai as genai
        from backend.config import settings

        genai.configure(api_key=settings.GEMINI_API_KEY)
        model = genai.GenerativeModel("gemini-2.0-flash")

        image_part = {"mime_type": "image/png", "data": img_bytes}
        prompt = (
            "Extract all text from this document page image exactly as it appears. "
            "Preserve headings, bullet points, tables, and paragraph structure. "
            "Return plain text only — no commentary, no markdown fences."
        )
        response = model.generate_content([prompt, image_part])
        return (response.text or "").strip()
    except Exception as e:
        logger.warning(f"Gemini Vision OCR failed for {filename} page {page_num}: {e}")
        return ""


def parse_document(file_path: str, filename: str) -> List[str]:
    """Parse a single uploaded document into source-tagged text chunks."""
    chunks: List[str] = []
    ext = os.path.splitext(filename)[1].lower().lstrip(".")

    try:
        # ── PDF ───────────────────────────────────────────────────────────
        if ext == "pdf":
            from pypdf import PdfReader

            reader = PdfReader(file_path)
            scanned_pages: List[int] = []

            for page_num, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ""
                # Strip NUL bytes immediately — PostgreSQL rejects them
                text = text.replace("\x00", "").strip()
                alpha_chars = sum(1 for c in text if c.isalpha())
                # Use pypdf text only if it is substantial (≥40 alpha chars AND
                # ≥300 total chars). Short extractions usually mean pypdf only
                # captured a footer/watermark and missed the real table content
                # (e.g. government PDFs, form-based admit cards).
                if text and alpha_chars >= 40 and len(text) >= 300:
                    tagged = f"[source: {filename} | page {page_num}]\n{text}"
                    chunks.append(tagged)
                else:
                    scanned_pages.append(page_num)

            # Fallback: use PyMuPDF to render scanned pages → Gemini Vision OCR
            if scanned_pages:
                logger.info(
                    f"{filename}: {len(scanned_pages)} scanned page(s) — running Gemini Vision OCR"
                )
                try:
                    import fitz  # PyMuPDF

                    doc = fitz.open(file_path)
                    for page_num in scanned_pages:
                        page = doc[page_num - 1]
                        pix = page.get_pixmap(dpi=200)
                        img_bytes = pix.tobytes("png")
                        ocr_text = _gemini_ocr_page(img_bytes, filename, page_num)
                        if ocr_text:
                            # Apply sliding window so long OCR pages get multiple chunks
                            windows = _sliding_window(ocr_text, size=500, overlap=80)
                            if not windows:
                                windows = [ocr_text]
                            for n, w in enumerate(windows, start=1):
                                suffix = f"part {n}" if len(windows) > 1 else "OCR"
                                tagged = (
                                    f"[source: {filename} | page {page_num} ({suffix})]\n{w}"
                                )
                                chunks.append(tagged)
                    doc.close()
                except ImportError:
                    logger.warning("PyMuPDF not installed — scanned PDF pages skipped.")
                except Exception as e:
                    logger.error(f"PyMuPDF/OCR error for {filename}: {e}")

        # ── PPTX ──────────────────────────────────────────────────────────
        elif ext == "pptx":
            from pptx import Presentation

            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    if not hasattr(shape, "text_frame"):
                        continue
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(
                            run.text for run in para.runs if run.text
                        ).strip()
                        if line:
                            texts.append(line)
                # Speaker notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        texts.append(f"[Speaker notes]: {notes}")

                joined = "\n".join(texts)
                if joined.strip():
                    tagged = f"[source: {filename} | slide {slide_num}]\n{joined}"
                    chunks.append(tagged)

        # ── CSV ───────────────────────────────────────────────────────────
        elif ext == "csv":
            import pandas as pd

            df = pd.read_csv(file_path, dtype=str)
            chunks.extend(_df_to_chunks(df, filename, "sheet1"))

        # ── XLSX / XLS ────────────────────────────────────────────────────
        elif ext in ("xlsx", "xls"):
            import pandas as pd

            xl = pd.ExcelFile(file_path)
            for sheet_name in xl.sheet_names:
                try:
                    df = pd.read_excel(xl, sheet_name=sheet_name, dtype=str)
                    sheet_label = f"sheet:{sheet_name}"
                    chunks.extend(_df_to_chunks(df, filename, sheet_label))
                except Exception as e:
                    logger.warning(f"Could not parse sheet '{sheet_name}' in {filename}: {e}")

        # ── DOCX ──────────────────────────────────────────────────────────
        elif ext == "docx":
            from docx import Document

            doc = Document(file_path)
            buf = []
            for para in doc.paragraphs:
                if para.text and para.text.strip():
                    buf.append(para.text.strip())
            for table in doc.tables:
                for row in table.rows:
                    cell_texts = [
                        cell.text.strip()
                        for cell in row.cells
                        if cell.text and cell.text.strip()
                    ]
                    if cell_texts:
                        buf.append(" | ".join(cell_texts))

            windows = _sliding_window("\n".join(buf))
            for n, w in enumerate(windows, start=1):
                chunks.append(f"[source: {filename} | chunk {n}]\n{w}")

        # ── TXT / MD ──────────────────────────────────────────────────────
        elif ext in ("txt", "md"):
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            windows = _sliding_window(text)
            for n, w in enumerate(windows, start=1):
                chunks.append(f"[source: {filename} | chunk {n}]\n{w}")

        else:
            logger.warning(f"Unsupported file extension: {ext}")

    except Exception as e:
        logger.error(f"Failed to parse {filename}: {e}", exc_info=True)

    return [c for c in chunks if c and c.strip()]
